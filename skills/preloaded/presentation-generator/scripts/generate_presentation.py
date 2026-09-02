#!/usr/bin/env python3
"""Generate matching PPTX and self-contained HTML presentation artifacts."""

from __future__ import annotations

import html
import json
import os
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_ENV = "WEKNORA_SKILL_OUTPUT_DIR"
NAME_RE = re.compile(r"[^A-Za-z0-9_-]+")


def read_payload() -> dict:
    raw = sys.stdin.read()
    if not raw.strip():
        raise ValueError("input JSON is required")
    value = json.loads(raw)
    if not isinstance(value, dict):
        raise ValueError("input must be a JSON object")
    return value


def normalise(payload: dict) -> tuple[str, str, str, str, list[dict]]:
    title = str(payload.get("title", "")).strip()
    if not title:
        raise ValueError("title is required")
    subtitle = str(payload.get("subtitle", "")).strip()
    author = str(payload.get("author", "")).strip()
    output_name = NAME_RE.sub("-", str(payload.get("output_name", "presentation"))).strip("-")
    output_name = output_name or "presentation"
    slides = payload.get("slides")
    if not isinstance(slides, list) or not slides:
        raise ValueError("slides must be a non-empty array")
    cleaned: list[dict] = []
    for item in slides[:30]:
        if not isinstance(item, dict):
            continue
        slide_title = str(item.get("title", "")).strip() or "Untitled"
        bullets = item.get("bullets", [])
        if not isinstance(bullets, list):
            bullets = []
        cleaned.append({"title": slide_title, "bullets": [str(x).strip() for x in bullets[:10] if str(x).strip()]})
    if not cleaned:
        raise ValueError("slides contains no valid slide object")
    return title, subtitle, author, output_name, cleaned


def add_title(slide, text: str, top: float, size: int, color: RGBColor) -> None:
    box = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(11.7), Inches(1.0))
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.LEFT
    paragraph.font.name = "Microsoft YaHei"
    paragraph.font.size = Pt(size)
    paragraph.font.bold = True
    paragraph.font.color.rgb = color


def build_pptx(title: str, subtitle: str, author: str, slides: list[dict], output: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    cover = prs.slides.add_slide(prs.slide_layouts[6])
    cover.background.fill.solid()
    cover.background.fill.fore_color.rgb = RGBColor(15, 36, 64)
    add_title(cover, title, 2.1, 34, RGBColor(255, 255, 255))
    if subtitle:
        add_title(cover, subtitle, 3.25, 20, RGBColor(109, 213, 170))
    if author:
        add_title(cover, author, 5.8, 12, RGBColor(190, 205, 218))

    for item in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(247, 250, 252)
        add_title(slide, item["title"], 0.55, 26, RGBColor(15, 36, 64))
        body = slide.shapes.add_textbox(Inches(1.0), Inches(1.75), Inches(11.3), Inches(4.9))
        frame = body.text_frame
        frame.clear()
        for index, bullet in enumerate(item["bullets"]):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
            paragraph.font.name = "Microsoft YaHei"
            paragraph.font.size = Pt(21)
            paragraph.font.color.rgb = RGBColor(38, 55, 69)
            paragraph.space_after = Pt(13)
            paragraph.text = "•  " + paragraph.text
        accent = slide.shapes.add_shape(1, Inches(0.55), Inches(0.55), Inches(0.08), Inches(5.95))
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(7, 192, 95)
        accent.line.fill.background()
    prs.save(output)


def build_html(title: str, subtitle: str, author: str, slides: list[dict], output: Path) -> None:
    pages = [{"title": title, "bullets": [x for x in (subtitle, author) if x]}] + slides
    sections = []
    for index, item in enumerate(pages):
        bullets = "".join(f"<li>{html.escape(value)}</li>" for value in item["bullets"])
        sections.append(
            f'<section class="slide"><span class="number">{index + 1}/{len(pages)}</span>'
            f'<h1>{html.escape(item["title"])}</h1><ul>{bullets}</ul></section>'
        )
    document = f"""<!doctype html><html lang="zh-CN"><head><meta charset="utf-8">
<meta name="viewport" content="width=device-width,initial-scale=1"><title>{html.escape(title)}</title>
<style>html,body{{margin:0;background:#0f2440;color:#fff;font-family:"Microsoft YaHei",sans-serif}}
.slide{{box-sizing:border-box;min-height:100vh;padding:9vh 8vw;display:none;background:linear-gradient(135deg,#0f2440,#183a5a)}}
.slide.active{{display:block}}h1{{font-size:clamp(32px,5vw,68px);margin:0 0 6vh;border-left:8px solid #07c05f;padding-left:24px}}
ul{{font-size:clamp(20px,2.3vw,34px);line-height:1.65}}.number{{position:absolute;right:3vw;bottom:3vh;color:#a9bdca}}
.hint{{position:fixed;left:3vw;bottom:3vh;color:#a9bdca;font-size:13px}}</style></head><body>
{''.join(sections)}<span class="hint">← → 切换</span><script>
const slides=[...document.querySelectorAll('.slide')];let current=0;
function show(n){{current=(n+slides.length)%slides.length;slides.forEach((s,i)=>s.classList.toggle('active',i===current));}}
addEventListener('keydown',e=>{{if(e.key==='ArrowRight'||e.key===' ')show(current+1);if(e.key==='ArrowLeft')show(current-1);}});show(0);
</script></body></html>"""
    output.write_text(document, encoding="utf-8")


def main() -> int:
    try:
        title, subtitle, author, output_name, slides = normalise(read_payload())
        output_dir = Path(os.environ.get(OUTPUT_ENV, "/workspace/output"))
        output_dir.mkdir(parents=True, exist_ok=True)
        pptx_path = output_dir / f"{output_name}.pptx"
        html_path = output_dir / f"{output_name}.html"
        build_pptx(title, subtitle, author, slides, pptx_path)
        build_html(title, subtitle, author, slides, html_path)
        print(json.dumps({
            "success": True,
            "artifacts": [
                {"path": pptx_path.name, "type": "presentation", "purpose": "download"},
                {"path": html_path.name, "type": "web", "purpose": "preview"},
            ],
        }, ensure_ascii=False))
        return 0
    except Exception as exc:  # noqa: BLE001 - CLI returns a structured failure
        print(json.dumps({"success": False, "error": str(exc)}, ensure_ascii=False), file=sys.stderr)
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
